"""Generate final project presentation (PPTX)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "期末專題簡報_校園美食地圖.pptx"

# Theme colors
NAVY = RGBColor(0x1A, 0x36, 0x5D)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_bar(slide, title: str, subtitle: str | None = None) -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)


def _add_bullets(slide, items: list[str], left=0.6, top=1.4, width=12.0, height=5.5, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)


def _add_code_block(slide, code: str, top=1.5, height=4.5):
    box = slide.shapes.add_shape(1, Inches(0.6), Inches(top), Inches(12.0), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
    box.line.color.rgb = RGBColor(0x55, 0x55, 0x55)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = code
    p.font.name = "Consolas"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)


def _add_table_slide(slide, headers: list[str], rows: list[list[str]], top=1.5):
    cols, row_count = len(headers), len(rows) + 1
    table_shape = slide.shapes.add_table(row_count, cols, Inches(0.6), Inches(top), Inches(12.0), Inches(0.45 * row_count))
    table = table_shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # --- Slide 1: Cover ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, NAVY)
    title_box = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.5))
    tf = title_box.text_frame
    tf.text = "校園附近美食地圖"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "國立宜蘭大學 · Python 程式期末專題"
    p2.font.size = Pt(22)
    p2.font.color.rgb = ORANGE
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = "作者：Boson Huang  |  雲端 Demo：https://niu-foodmap.streamlit.app"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
    p3.alignment = PP_ALIGN.CENTER

    # --- Slide 2: Agenda ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, LIGHT_BG)
    _add_title_bar(s, "簡報大綱")
    _add_bullets(
        s,
        [
            "(1) 動機 — 問題與任務描述",
            "(2) 目標 — 欲達成之目標",
            "(3) 軟硬體工具與執行平台",
            "(4) 組員分工事項",
            "(5) 架構設計與實作過程",
            "(6) 遇到的困難與解決方式",
        ],
        size=20,
    )

    # --- Slide 3: Motivation ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, LIGHT_BG)
    _add_title_bar(s, "(1) 動機", "問題或任務描述")
    _add_bullets(
        s,
        [
            "【現況痛點】宜大學生在校外覓食時，難以在「步行／單車可達」範圍內，",
            "　依網路評價（星等＋評論量）快速排序、篩選餐廳。",
            "",
            "【具體問題】",
            "• Google Maps 預設排序 ≠ 距離 + 評價綜合考量",
            "• 評論極少卻五星的店家容易誤導（需貝氏／加權修正）",
            "• 遠距高分店 vs 近距普通店 — 難以一次比較",
            "• 選擇困難：不知道「今天吃什麼」",
            "",
            "【任務】開發「宜大校本部周邊美食地圖」系統，",
            "　整合距離計算、多種評分算法、互動地圖與轉盤選店功能。",
        ],
        size=17,
    )

    # --- Slide 4: Goals ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, LIGHT_BG)
    _add_title_bar(s, "(2) 目標", "欲達成之目標")
    _add_bullets(
        s,
        [
            "【功能目標】",
            "• FR-1：依校園中心座標 + 半徑（km）篩選周邊餐廳",
            "• FR-2：四種排序 — 綜合分／貝氏星等／黃氏星等／距離",
            "• FR-3：CLI 輸出 table／json，含 Google Maps 連結",
            "• FR-4：Streamlit 雙分頁 — 美食地圖 + 轉盤選擇器",
            "• FR-5：支援自訂 JSON 資料（Mock 或 Google Places 快取）",
            "",
            "【品質目標】",
            "• 67 項 pytest 全綠，coverage ≥ 70%",
            "• 查詢時不打 API — 離線快取，p95 延遲 ≤ 120ms",
            "• 部署至 Streamlit Cloud，供系上／同學即時使用",
        ],
        size=17,
    )

    # --- Slide 5: Tools ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, LIGHT_BG)
    _add_title_bar(s, "(3) 軟硬體工具與執行平台")
    _add_table_slide(
        s,
        ["類別", "工具／平台", "用途"],
        [
            ["語言", "Python 3.10+", "核心邏輯、CLI、測試"],
            ["Web UI", "Streamlit 1.28+", "互動式網頁介面"],
            ["地圖", "PyDeck + Carto Voyager", "點位地圖視覺化"],
            ["資料", "Pandas", "表格展示與篩選"],
            ["資料來源", "Google Places API (New)", "一次性抓取 → JSON 快取"],
            ["測試", "pytest + pytest-cov", "67 tests，CI gate"],
            ["CI/CD", "GitHub Actions", "自動測試 + benchmark gate"],
            ["部署", "Streamlit Community Cloud", "https://niu-foodmap.streamlit.app"],
            ["版本控制", "Git + GitHub (boson316/niu-foodmap)", "協作與自動部署"],
            ["硬體", "Windows PC / 一般筆電", "本機開發與 demo"],
            ["瀏覽器", "Chrome / Edge / 手機瀏覽器", "雲端 app 存取"],
        ],
        top=1.35,
    )

    # --- Slide 6: Team ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, LIGHT_BG)
    _add_title_bar(s, "(4) 組員分工事項", "※ 若為小組作業，請依實際組員修改姓名欄")
    _add_table_slide(
        s,
        ["成員", "負責模組", "主要工作"],
        [
            ["Boson Huang（主開發）", "scoring.py / service.py", "黃氏星等、貝氏星等、綜合分算法；排序服務層"],
            ["【組員 A】", "providers.py / scripts/", "Google Places 格網抓取、JSON 快取管線"],
            ["【組員 B】", "streamlit_app.py", "Streamlit UI、地圖互動、側欄篩選"],
            ["【組員 C】", "wheel.py / visit_counter.py", "轉盤 SVG、瀏覽人次計數"],
            ["全員", "tests/ + CI", "TDD 測試撰寫、GitHub Actions 維護、文件"],
        ],
        top=1.5,
    )
    note = s.shapes.add_textbox(Inches(0.6), Inches(5.8), Inches(12), Inches(1))
    note.text_frame.text = "※ 本 repo 主要作者為 Boson Huang；若為個人專題，以上分工可合併為單人全棧開發。"
    note.text_frame.paragraphs[0].font.size = Pt(12)
    note.text_frame.paragraphs[0].font.color.rgb = GRAY

    # --- Slide 7: Architecture ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, LIGHT_BG)
    _add_title_bar(s, "(5) 架構設計", "系統分層與資料流")
    flow = (
        "[Google Places API] ──fetch 一次──► data/places_cache.json\n"
        "[places_cache.public.json] ──雲端──┤\n"
        "[sample_restaurants.json] ──Mock───┤\n"
        "                                    ▼\n"
        "                         MockRestaurantProvider\n"
        "                                    ▼\n"
        "                         FoodMapService.rank_nearby()\n"
        "                         (bbox → haversine → heapq top-k)\n"
        "                                    ▼\n"
        "              ┌─────────────────────┴─────────────────────┐\n"
        "              ▼                                           ▼\n"
        "       CLI (table/json)                    Streamlit（雙分頁）\n"
        "                                           ├─ 🗺️ 美食地圖\n"
        "                                           └─ 🎡 轉盤選擇"
    )
    _add_code_block(s, flow, top=1.35, height=5.5)

    # --- Slide 8: Scoring ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, LIGHT_BG)
    _add_title_bar(s, "(5) 核心算法", "黃氏星等 + 綜合分數")
    code = (
        "# 黃氏星等（scoring.py）\n"
        "huang = rating × 分級係數 × (0.5 + 0.5×rating/5)\n"
        "                  × (0.4 + 0.6×log1p(評論數)/log1p(500))\n"
        "# 分級：1~2星×0.65 | 3星×1.0 | 4~5星×1.25\n\n"
        "# 綜合分數（預設排序 & 轉盤候選池 Top 30）\n"
        "composite = huang_rating × exp(-distance_km / decay_km)\n"
        "# decay_km 預設 0.6 → 距離愈近分數愈高\n\n"
        "# 貝氏星等（抑制少評論極端值）\n"
        "bayes = (prior×prior_mean + rating×count) / (prior + count)"
    )
    _add_code_block(s, code, top=1.35, height=4.8)
    _add_bullets(
        s,
        ["轉盤：綜合分 Top 30 等機率抽選  |  地圖：PyDeck 點狀標記 + zoom 15/16"],
        top=6.3,
        height=0.8,
        size=14,
    )

    # --- Slide 9: Implementation steps ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, LIGHT_BG)
    _add_title_bar(s, "(5) 實作過程與方法")
    _add_bullets(
        s,
        [
            "Step 1 — 需求分析：PRD 定義 FR-1~5、宜大校本部座標 (24.7464, 121.7457)",
            "Step 2 — TDD 開發：先寫 tests/，再實作 scoring → distance → service → CLI",
            "Step 3 — 資料管線：Google Places 格網密集抓取 (~100 家) → JSON 離線快取",
            "Step 4 — Streamlit UI：表格 11 欄 + 點列連動地圖 + 側欄篩選",
            "Step 5 — 轉盤功能：SVG 扇形盤 + 角度對齊 helper（Python/JS 同邏輯可測）",
            "Step 6 — 品質保證：67 tests、integrity SHA256 驗證、GitHub Actions CI",
            "Step 7 — 雲端部署：push → Streamlit Cloud 自動 redeploy",
        ],
        size=17,
    )

    # --- Slide 10: Streamlit UI ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, LIGHT_BG)
    _add_title_bar(s, "(5) Streamlit 功能一覽")
    _add_table_slide(
        s,
        ["分頁", "功能", "技術細節"],
        [
            ["🗺️ 美食地圖", "餐廳清單 + 互動地圖", "st.dataframe 11 欄；點列 → 地圖定位"],
            ["", "Google Maps 一鍵跳轉", "LinkColumn + st.link_button"],
            ["", "綜合分數說明", "公式 + 實例（Muze／炸醬麵）"],
            ["🎡 轉盤選擇", "隨機選店", "綜合分 Top 30 等機率抽選"],
            ["", "SVG 扇形轉盤", "指針 = 結果 # 一致；多次旋轉不偏移"],
            ["頁尾", "作者 + 瀏覽人次", "counterapi.dev 瀏覽器端計數"],
        ],
        top=1.35,
    )

    # --- Slide 11: Service code ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, LIGHT_BG)
    _add_title_bar(s, "(5) 程式碼節錄", "service.py — 效能路徑")
    code = (
        "# FoodMapService.rank_nearby() 核心流程\n"
        "1. bbox 粗篩（approx_in_radius_bbox）\n"
        "2. haversine 精算距離（distance.py）\n"
        "3. 計算 bayesian / huang / composite 分數\n"
        "4. heapq 取 top-k（O(n log k)，非全排序）\n\n"
        "def build_maps_url(name, lat, lon, place_id=None):\n"
        "    if place_id:\n"
        "        return f'...query={name}&query_place_id={pid}'\n"
        "    return f'...query={name}@{lat},{lon}'"
    )
    _add_code_block(s, code, top=1.35, height=5.2)

    # --- Slide 12: Difficulties ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, LIGHT_BG)
    _add_title_bar(s, "(6) 遇到的困難與解決方式")
    _add_table_slide(
        s,
        ["困難", "原因", "解決方式"],
        [
            ["轉盤指針 ≠ 結果 #", "CSS conic-gradient 與標籤層偏移 ~7 格", "改 SVG 扇形 + spin_delta_degrees 增量旋轉"],
            ["LinkColumn 顯示「店名文字」", "Streamlit 1.55 中文 display_text bug", "店名 TextColumn + 點列選店連動地圖"],
            ["雲端瀏覽人次破圖", "SeeYouFarm 404；伺服器端 urllib 被擋", "改瀏覽器端 counterapi.dev + sessionStorage 防灌票"],
            ["CI coverage 68% 紅燈", "streamlit_app.py 計入 coverage", ".coveragerc omit UI 層；核心 ≥70%"],
            ["Cloud Python 3.14 怪錯", "dataclass import race", ".python-version 釘 3.10.12"],
            ["integrity hash 失敗", "Windows CRLF vs Linux LF", "LF 正規化 + 更新 SHA256 hash"],
            ["Mock 資料與真實落差", "15 筆範例不足 demo", "Google Places 格網抓取 ~100 家快取"],
        ],
        top=1.3,
    )

    # --- Slide 13: Demo ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, LIGHT_BG)
    _add_title_bar(s, "Demo 與成果")
    _add_bullets(
        s,
        [
            "【雲端 Demo】https://niu-foodmap.streamlit.app",
            "【GitHub】https://github.com/boson316/niu-foodmap",
            "",
            "【本機執行】",
            "  cd python程式期末專案_美食系統",
            "  pip install -r requirements.txt",
            "  $env:PYTHONPATH = \"src\"",
            "  streamlit run src/streamlit_app.py",
            "",
            "【測試】pytest -q --cov=src --cov-fail-under=70  →  67 passed, ~86% coverage",
            "",
            "【成果】宜大周邊 ~100 家餐廳 · 四種排序 · 互動地圖 · 轉盤選店 · CI 全綠",
        ],
        size=18,
    )

    # --- Slide 14: Q&A ---
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, NAVY)
    box = s.shapes.add_textbox(Inches(0), Inches(2.8), Inches(13.33), Inches(2))
    tf = box.text_frame
    tf.text = "Q & A"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "謝謝聆聽"
    p2.font.size = Pt(24)
    p2.font.color.rgb = ORANGE
    p2.alignment = PP_ALIGN.CENTER

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Generated: {path}")
